#!/usr/bin/env python3
"""
Generate PowerPoint presentation for EMS Marikina application.
Run: python3 generate_presentation.py
Output: EMS_MARIKINA_PRESENTATION.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Color scheme
PRIMARY_COLOR = RGBColor(25, 56, 109)      # Dark blue
ACCENT_COLOR = RGBColor(220, 53, 69)       # Red
SUCCESS_COLOR = RGBColor(40, 167, 69)      # Green
LIGHT_GRAY = RGBColor(248, 249, 250)       # Light gray
DARK_TEXT = RGBColor(33, 37, 41)           # Dark gray

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = ACCENT_COLOR
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    # Title text
    title_frame = title_shape.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    title_frame = title_shape.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Left column header
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(0.4))
    left_h = left_header.text_frame.paragraphs[0]
    left_h.text = left_title
    left_h.font.size = Pt(20)
    left_h.font.bold = True
    left_h.font.color.rgb = ACCENT_COLOR
    
    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(4.2), Inches(5.3))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_items):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)
    
    # Right column header
    right_header = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(4), Inches(0.4))
    right_h = right_header.text_frame.paragraphs[0]
    right_h.text = right_title
    right_h.font.size = Pt(20)
    right_h.font.bold = True
    right_h.font.color.rgb = ACCENT_COLOR
    
    # Right column content
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.7), Inches(4.0), Inches(5.3))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_items):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)
    
    return slide

# Slide 1: Title Slide
add_title_slide(prs, "EMS MARIKINA", "Emergency Medical Services Dispatch System")

# Slide 2: Introduction/App Description
add_content_slide(prs, "Introduction", [
    "✓ A real-time emergency dispatch system for Marikina Government",
    "✓ Coordinates ambulances, dispatchers, operators, and hospitals",
    "✓ Uses AI to recommend optimal hospital based on distance & capacity",
    "✓ Reduces emergency response time significantly",
    "✓ Provides complete audit trail for quality assurance",
    "✓ Multi-role system (Admin, Dispatcher, Operator, Driver)"
])

# Slide 3: Problem Statement
add_content_slide(prs, "Problem Statement", [
    "Current System Challenges:",
    "• Long response times from emergency call to ambulance dispatch",
    "• Manual hospital selection leads to inefficient routing",
    "• No real-time coordination between stakeholders",
    "• Poor tracking of incident details and outcomes",
    "• Difficulty in emergency response optimization"
])

# Slide 4: Key Features (Part 1)
add_content_slide(prs, "Key Features - Part 1", [
    "🔐 Authentication & Authorization",
    "    - Email/password login with BCrypt encryption",
    "    - Role-based access control & permissions",
    "",
    "📋 Dispatch Management",
    "    - Create emergency dispatches with patient details",
    "    - AI recommends optimal hospital automatically",
    "    - One-tap driver assignment"
])

# Slide 5: Key Features (Part 2)
add_content_slide(prs, "Key Features - Part 2", [
    "📍 Real-Time Monitoring",
    "    - Live driver location tracking (GPS)",
    "    - Dispatch board with all active incidents",
    "    - Instant SignalR notifications",
    "",
    "📊 Analytics & Audit",
    "    - Complete dispatch history",
    "    - Session logs for security",
    "    - Performance metrics dashboard"
])

# Slide 6: User Roles
add_two_column_slide(prs, "User Roles & Responsibilities",
    "Dispatcher & Admin",
    [
        "Dispatcher:",
        "- Receives emergency calls",
        "- Creates dispatches",
        "- Assigns drivers",
        "- Views recommendations",
        "",
        "Admin:",
        "- System oversight",
        "- User management",
        "- Analytics viewing"
    ],
    "EMS Operator & Driver",
    [
        "EMS Operator:",
        "- Monitors live operations",
        "- Coordinates drivers",
        "- Views driver locations",
        "",
        "Driver:",
        "- Receives dispatch assignment",
        "- Navigates to patient",
        "- Updates dispatch status"
    ]
)

# Slide 7: Application Workflow
add_content_slide(prs, "Application Workflow", [
    "1️⃣ CALL RECEIVED → Dispatcher records emergency call",
    "",
    "2️⃣ DISPATCH CREATED → Patient info & location entered",
    "",
    "3️⃣ AI RECOMMENDATION → System recommends best hospital",
    "",
    "4️⃣ DRIVER ASSIGNED → One-tap assignment to driver",
    "",
    "5️⃣ REAL-TIME TRACKING → Operator monitors progress",
    "",
    "6️⃣ HOSPITAL HANDOVER → Patient delivered, dispatch completed"
])

# Slide 8: HCI Principles - Part 1
add_content_slide(prs, "HCI Principles Applied - Part 1", [
    "🎯 Know the User",
    "    Why: Different roles have different needs and stress levels",
    "    How: Role-specific dashboards, quick-access high-urgency buttons",
    "",
    "👥 User-Centered Design",
    "    Why: Users interact during high-stress emergencies",
    "    How: Large buttons, minimal taps, color-coded severity, real-time feedback"
])

# Slide 9: HCI Principles - Part 2
add_content_slide(prs, "HCI Principles Applied - Part 2", [
    "🔄 Consistency Across Interfaces",
    "    Why: Users switch between roles frequently",
    "    How: Same colors, typography, buttons across all dashboards",
    "",
    "⚠️ Error Prevention & Recovery",
    "    Why: Emergency mistakes have serious consequences",
    "    How: Confirmation dialogs, field validation, auto-save functionality"
])

# Slide 10: Usability Principles - Part 1
add_two_column_slide(prs, "Usability Principles Applied",
    "Learnability & Efficiency",
    [
        "Learnability:",
        "- Intuitive icons & labels",
        "- Consistent layout everywhere",
        "- Clear visual hierarchy",
        "",
        "Efficiency:",
        "- Keyboard shortcuts",
        "- Auto-complete fields",
        "- One-tap common actions",
        "- Data caching"
    ],
    "Memorability & Error Tolerance",
    [
        "Memorability:",
        "- Consistent commands",
        "- Familiar icons",
        "- Standard interactions",
        "",
        "Error Tolerance:",
        "- Undo functionality",
        "- Clear error messages",
        "- Graceful degradation",
        "- Auto-recovery"
    ]
)

# Slide 11: Usability Principles - Part 2
add_content_slide(prs, "Usability Principles - Accessibility & Responsiveness", [
    "♿ Accessibility",
    "    - High contrast mode support",
    "    - Configurable font sizes",
    "    - Screen reader compatibility",
    "    - Keyboard navigation",
    "",
    "📱 Responsiveness",
    "    - Mobile-first design (phones, tablets, desktops)",
    "    - Offline mode support",
    "    - Progressive loading",
    "    - Works on low-bandwidth connections"
])

# Slide 12: Technology Stack
add_two_column_slide(prs, "Technology Stack",
    "Frontend & Backend",
    [
        "Frontend:",
        "- Flutter (Web)",
        "- Riverpod (State)",
        "- SignalR (Real-time)",
        "- GoRouter (Navigation)",
        "",
        "Backend:",
        "- .NET 8",
        "- ASP.NET Core",
        "- JWT Authentication",
        "- SignalR Core"
    ],
    "Database & Deployment",
    [
        "Database:",
        "- MongoDB",
        "- Indexed collections",
        "- TTL session cleanup",
        "",
        "Deployment:",
        "- Docker containers",
        "- Docker Compose",
        "- CI/CD pipeline",
        "- 99.9% uptime SLA"
    ]
)

# Slide 13: Database Schema
add_content_slide(prs, "Database Collections", [
    "📄 Users: Stores all user accounts with encrypted passwords",
    "",
    "🚨 Dispatches: Emergency incidents with patient, location, status, hospital recommendations",
    "",
    "🚗 Drivers: Driver profiles linked to users with current location & status",
    "",
    "🔐 Sessions: Login sessions with TTL auto-expiration (24 hours)",
    "",
    "🏥 EMSBases: EMS base stations and hospital locations"
])

# Slide 14: Security Features
add_content_slide(prs, "Security & Performance", [
    "🔒 Security Features:",
    "    - BCrypt password hashing (salt rounds: 12)",
    "    - JWT token authentication (expiry: 8 hours)",
    "    - Session TTL with auto-cleanup (24 hours)",
    "    - Parameterized MongoDB queries (SQL injection prevention)",
    "    - HTTPS/TLS encryption",
    "",
    "⚡ Performance Metrics:",
    "    - Dispatch creation: < 2 seconds",
    "    - Driver assignment: < 500ms",
    "    - Real-time latency: < 1 second",
    "    - Concurrent support: 100+ drivers"
])

# Slide 15: API Overview
add_content_slide(prs, "API Endpoints Overview", [
    "🔑 Authentication:",
    "    POST /api/auth/login, POST /api/auth/logout",
    "",
    "📋 Dispatches:",
    "    GET /api/dispatches, GET /api/dispatches/my (Driver only)",
    "    POST /api/dispatches, PUT /api/dispatches/{id}/assign",
    "",
    "🚗 Drivers:",
    "    GET /api/drivers, PUT /api/drivers/location",
    "",
    "👥 Users:",
    "    GET /api/users, POST /api/users, PUT /api/users/{id}"
])

# Slide 16: Installation & Testing
add_content_slide(prs, "Installation & Testing Credentials", [
    "Quick Start:",
    "1. docker compose up (starts backend & database)",
    "2. flutter run -d web (starts frontend)",
    "3. Access at http://localhost:64262",
    "",
    "Test Accounts:",
    "• Admin: admin@ems.local / Admin@123",
    "• Dispatcher: dispatcher@ems.local / Dispatch@123",
    "• EMS Operator: operator@ems.local / Operator@123",
    "• Driver: driver@ems.local / Driver@123"
])

# Slide 17: Future Enhancements
add_content_slide(prs, "Future Enhancements & Roadmap", [
    "🔄 Phase 2 Improvements:",
    "    - Native iOS/Android mobile apps",
    "    - Advanced analytics dashboard",
    "    - Hospital EMR system integration",
    "    - SMS notification fallback",
    "",
    "🌟 Phase 3+ Features:",
    "    - Multi-language support (Tagalog/English)",
    "    - Voice commands for hands-free dispatch",
    "    - AR navigation for drivers",
    "    - Predictive analytics (anticipate emergencies)"
])

# Slide 18: Conclusion
add_title_slide(prs, "Thank You", "EMS Marikina - Saving Lives Through Technology")

# Save presentation
output_file = "/vercel/share/v0-project/EMS_MARIKINA_PRESENTATION.pptx"
prs.save(output_file)
print(f"✅ Presentation created successfully: {output_file}")
